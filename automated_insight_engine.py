"""
Automated Insight Engine
Analyzes marketing campaign data and generates PowerPoint reports with AI insights.
"""

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import os
import sys
import google.generativeai as genai

# Set console encoding to UTF-8
if sys.platform == 'win32':
    import codecs
    sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')

# Configuration
GEMINI_API_KEY = "AIzaSyBprEGKV91oF8DHIf3LB86LqGT-CRlPJ0E"
CSV_FILE = "marketing_campaign.csv"
CTR_CHART = "ctr_chart.png"
CPA_CHART = "cpa_chart.png"
OUTPUT_PPT = "Campaign_Report.pptx"

print("=" * 60)
print("AUTOMATED INSIGHT ENGINE")
print("=" * 60)

# ============================================================================
# PHASE 1: DATA INGESTION & VISUALIZATION
# ============================================================================
print("\n[PHASE 1] Data Ingestion & Visualization")
print("-" * 60)

# Load data
print(f"Loading {CSV_FILE}...")
df = pd.read_csv(CSV_FILE, sep=';')
print(f"[OK] Loaded {len(df)} customer records with {len(df.columns)} columns")

# Clean data - handle missing Income values
print("\nCleaning data...")
missing_income = df['Income'].isnull().sum()
if missing_income > 0:
    median_income = df['Income'].median()
    df.loc[df['Income'].isnull(), 'Income'] = median_income
    print(f"[OK] Filled {missing_income} missing Income values with median: ${median_income:,.0f}")

# Calculate total spending per customer
df['TotalSpend'] = (df['MntWines'] + df['MntFruits'] + df['MntMeatProducts'] + 
                     df['MntFishProducts'] + df['MntSweetProducts'] + df['MntGoldProds'])

# Transform customer data into campaign metrics
print("\nTransforming customer data into campaign metrics...")

# DYNAMIC CAMPAIGN DETECTION - Auto-detect all campaign columns
# Looks for columns like: AcceptedCmp1, AcceptedCmp2, ..., AcceptedCmpN, Response
campaigns = {}

# Find all AcceptedCmp columns
cmp_columns = [col for col in df.columns if col.startswith('AcceptedCmp')]
# Sort them numerically (AcceptedCmp1, AcceptedCmp2, etc.)
cmp_columns_sorted = sorted(cmp_columns, key=lambda x: int(x.replace('AcceptedCmp', '')))

for col in cmp_columns_sorted:
    campaign_num = col.replace('AcceptedCmp', '')
    campaigns[f'Campaign {campaign_num}'] = col

# Add Response campaign if it exists
if 'Response' in df.columns:
    campaigns['Response Campaign'] = 'Response'

print(f"Auto-detected {len(campaigns)} campaigns: {list(campaigns.keys())}")

campaign_metrics = []
total_customers = len(df)

for campaign_name, column in campaigns.items():
    # Count acceptances (conversions)
    acceptances = df[column].sum()
    
    # Calculate acceptance rate (CTR equivalent)
    acceptance_rate = (acceptances / total_customers) * 100
    
    # Calculate average spend per acceptor (CPA equivalent - but as revenue metric)
    # For customers who accepted, what's their average total spend
    acceptors = df[df[column] == 1]
    if len(acceptors) > 0:
        avg_spend_per_acceptor = acceptors['TotalSpend'].mean()
    else:
        avg_spend_per_acceptor = 0
    
    campaign_metrics.append({
        'Campaign': campaign_name,
        'Acceptances': int(acceptances),
        'AcceptanceRate': acceptance_rate,
        'AvgSpendPerAcceptor': avg_spend_per_acceptor
    })
    
    print(f"  {campaign_name}: {acceptances} acceptances ({acceptance_rate:.2f}%), Avg Spend: ${avg_spend_per_acceptor:.2f}")

# Create metrics DataFrame
metrics_df = pd.DataFrame(campaign_metrics)

# Generate Chart 1: Campaign Acceptance Rate (CTR)
print("\nGenerating CTR Chart...")
plt.figure(figsize=(12, 6))
colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A', '#98D8C8', '#F7DC6F']
bars = plt.bar(metrics_df['Campaign'], metrics_df['AcceptanceRate'], color=colors, edgecolor='black', linewidth=1.2)

# Add value labels on bars
for bar in bars:
    height = bar.get_height()
    plt.text(bar.get_x() + bar.get_width()/2., height,
             f'{height:.1f}%',
             ha='center', va='bottom', fontsize=10, fontweight='bold')

plt.xlabel('Campaign Name', fontsize=12, fontweight='bold')
plt.ylabel('Acceptance Rate (%)', fontsize=12, fontweight='bold')
plt.title('Campaign Engagement Analysis - Acceptance Rates', fontsize=14, fontweight='bold', pad=20)
plt.xticks(rotation=45, ha='right')
plt.grid(axis='y', alpha=0.3, linestyle='--')
plt.tight_layout()
plt.savefig(CTR_CHART, dpi=300, bbox_inches='tight')
print(f"[OK] Saved {CTR_CHART}")
plt.close()

# Generate Chart 2: Average Spend per Acceptor (ROI metric)
print("Generating CPA Chart...")
plt.figure(figsize=(12, 6))
bars = plt.bar(metrics_df['Campaign'], metrics_df['AvgSpendPerAcceptor'], color=colors, edgecolor='black', linewidth=1.2)

# Add value labels on bars
for bar in bars:
    height = bar.get_height()
    plt.text(bar.get_x() + bar.get_width()/2., height,
             f'${height:.0f}',
             ha='center', va='bottom', fontsize=10, fontweight='bold')

plt.xlabel('Campaign Name', fontsize=12, fontweight='bold')
plt.ylabel('Average Spend per Acceptor ($)', fontsize=12, fontweight='bold')
plt.title('Campaign ROI Analysis - Average Customer Spend', fontsize=14, fontweight='bold', pad=20)
plt.xticks(rotation=45, ha='right')
plt.grid(axis='y', alpha=0.3, linestyle='--')
plt.tight_layout()
plt.savefig(CPA_CHART, dpi=300, bbox_inches='tight')
print(f"[OK] Saved {CPA_CHART}")
plt.close()

print("\n[OK] Phase 1 Complete: Charts generated successfully")

# ============================================================================
# PHASE 2: AI INSIGHT GENERATION
# ============================================================================
print("\n[PHASE 2] AI Insight Generation")
print("-" * 60)

# Prepare data summary for LLM
data_summary = f"""
Total Customers Analyzed: {total_customers:,}

Campaign Performance Metrics:
"""

for _, row in metrics_df.iterrows():
    data_summary += f"\n{row['Campaign']}:"
    data_summary += f"\n  - Acceptances: {row['Acceptances']}"
    data_summary += f"\n  - Acceptance Rate: {row['AcceptanceRate']:.2f}%"
    data_summary += f"\n  - Avg Spend per Acceptor: ${row['AvgSpendPerAcceptor']:.2f}"

# Construct the LLM prompt
prompt = f"""Act as a Senior Marketing Strategist. Analyze the following campaign performance data:

{data_summary}

Please provide:

1. A 3-bullet point executive summary of the overall campaign performance.

2. Identify the best performing campaign based on Acceptance Rate (highest engagement).

3. Identify the most valuable campaign based on Average Spend per Acceptor (highest customer value).

4. Recommend which campaign to optimize or reduce budget for based on the metrics.

Keep the tone professional and concise. Format your response clearly with headers for each section.
"""

# Call Gemini API
print("Calling Gemini API for insights...")
try:
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel('gemini-2.0-flash-exp')
    response = model.generate_content(prompt)
    ai_insights = response.text
    print("[OK] AI insights generated successfully")
    print("\nAI Insights Preview:")
    print("-" * 60)
    print(ai_insights[:300] + "..." if len(ai_insights) > 300 else ai_insights)
except Exception as e:
    print(f"[WARNING] API Error: {e}")
    print("Using fallback insights...")
    
    # Fallback insights if API fails
    best_ctr = metrics_df.loc[metrics_df['AcceptanceRate'].idxmax()]
    best_spend = metrics_df.loc[metrics_df['AvgSpendPerAcceptor'].idxmax()]
    worst_metric = metrics_df.loc[metrics_df['AcceptanceRate'].idxmin()]
    
    ai_insights = f"""EXECUTIVE SUMMARY

* Overall campaign performance shows varied engagement, with acceptance rates ranging from {metrics_df['AcceptanceRate'].min():.1f}% to {metrics_df['AcceptanceRate'].max():.1f}%.

* Total campaign acceptances: {metrics_df['Acceptances'].sum():,} customers across all campaigns.

* Customer spending patterns indicate strong value in accepted campaigns, with average spends ranging from ${metrics_df['AvgSpendPerAcceptor'].min():.0f} to ${metrics_df['AvgSpendPerAcceptor'].max():.0f}.

BEST PERFORMING CAMPAIGN (Engagement)

{best_ctr['Campaign']} achieved the highest acceptance rate at {best_ctr['AcceptanceRate']:.2f}%, demonstrating superior customer engagement and targeting effectiveness.

MOST VALUABLE CAMPAIGN (Revenue)

{best_spend['Campaign']} generated the highest average spend per acceptor at ${best_spend['AvgSpendPerAcceptor']:.0f}, indicating strong customer value and purchasing power.

RECOMMENDATION

Consider optimizing {worst_metric['Campaign']}, which showed the lowest acceptance rate at {worst_metric['AcceptanceRate']:.2f}%. Re-evaluate targeting criteria, messaging, or budget allocation to improve performance or redirect resources to higher-performing campaigns.
"""

print("\n[OK] Phase 2 Complete: Insights generated")

# ============================================================================
# PHASE 3: POWERPOINT REPORT ASSEMBLY
# ============================================================================
print("\n[PHASE 3] PowerPoint Report Assembly")
print("-" * 60)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
print("Creating Slide 1: Title...")
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Add title
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Marketing Performance Report"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.name = 'Calibri'
title_para.alignment = 1  # Center

# Add date
date_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
date_frame = date_box.text_frame
date_frame.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
date_para = date_frame.paragraphs[0]
date_para.font.size = Pt(20)
date_para.font.name = 'Calibri'
date_para.alignment = 1  # Center

print("[OK] Slide 1 created")

# Slide 2: Engagement Analysis
print("Creating Slide 2: Engagement Analysis...")
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Add title
title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Engagement Analysis"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.name = 'Calibri'

# Insert CTR chart
slide2.shapes.add_picture(CTR_CHART, Inches(0.75), Inches(1.2), width=Inches(8.5))
print("[OK] Slide 2 created with CTR chart")

# Slide 3: ROI & Recommendations
print("Creating Slide 3: ROI & Recommendations...")
slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Add title
title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "ROI & Recommendations"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.name = 'Calibri'

# Insert CPA chart (smaller to make room for text)
slide3.shapes.add_picture(CPA_CHART, Inches(0.5), Inches(1.1), width=Inches(5.5))

# Add AI insights text box
insights_box = slide3.shapes.add_textbox(Inches(6.2), Inches(1.1), Inches(3.3), Inches(5.5))
insights_frame = insights_box.text_frame
insights_frame.text = ai_insights
insights_frame.word_wrap = True

# Format insights text
for paragraph in insights_frame.paragraphs:
    paragraph.font.size = Pt(9)
    paragraph.font.name = 'Calibri'
    paragraph.space_after = Pt(6)

print("[OK] Slide 3 created with CPA chart and AI insights")

# Save presentation
prs.save(OUTPUT_PPT)
print(f"\n[OK] Presentation saved: {OUTPUT_PPT}")

print("\n[OK] Phase 3 Complete: PowerPoint report assembled")

# ============================================================================
# SUMMARY
# ============================================================================
print("\n" + "=" * 60)
print("AUTOMATED INSIGHT ENGINE - COMPLETE")
print("=" * 60)
print(f"\nGenerated Files:")
print(f"  1. {CTR_CHART} - Campaign engagement chart")
print(f"  2. {CPA_CHART} - Campaign ROI chart")
print(f"  3. {OUTPUT_PPT} - Complete PowerPoint report")
print(f"\nTotal Campaigns Analyzed: {len(campaigns)}")
print(f"Total Customers: {total_customers:,}")
print(f"Report Date: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
print("\n[OK] All tasks completed successfully!")
print("=" * 60)
