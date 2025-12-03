"""
Universal Insight Engine - Handles ANY CSV Type
Automatically detects: Marketing Campaigns, E-commerce Sales, or Generic Data
"""

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import os
import tempfile
import google.generativeai as genai
from io import BytesIO
import csv

# Configuration
GEMINI_API_KEY = "AIzaSyBprEGKV91oF8DHIf3LB86LqGT-CRlPJ0E"

# Page config
st.set_page_config(
    page_title="Universal Insight Engine",
    page_icon="📊",
    layout="wide"
)

# Custom CSS
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: bold;
        background: linear-gradient(90deg, #1f77b4, #2ca02c);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-align: center;
        margin-bottom: 1rem;
    }
    .sub-header {
        font-size: 1.2rem;
        color: #666;
        text-align: center;
        margin-bottom: 2rem;
    }
    .insight-box {
        background-color: #f0f8ff;
        padding: 1.5rem;
        border-radius: 10px;
        border-left: 5px solid #1f77b4;
        margin: 1rem 0;
    }
    .success-box {
        background-color: #d4f4dd;
        padding: 1rem;
        border-radius: 5px;
        border-left: 5px solid #28a745;
    }
    .dataset-badge {
        display: inline-block;
        padding: 0.5rem 1rem;
        border-radius: 20px;
        font-weight: bold;
        margin: 0.5rem;
    }
</style>
""", unsafe_allow_html=True)

# Helper Functions
def detect_delimiter(file_content):
    """Auto-detect CSV delimiter"""
    try:
        sample = file_content[:1024].decode('utf-8', errors='ignore')
        sniffer = csv.Sniffer()
        delimiter = sniffer.sniff(sample).delimiter
        return delimiter
    except:
        for delim in [',', ';', '\t', '|']:
            if delim in sample:
                return delim
        return ','

def detect_dataset_type(df):
    """Detect if dataset is Campaign, Sales, or Generic"""
    cols_lower = [col.lower() for col in df.columns]
    
    # Check for e-commerce/sales indicators
    sales_indicators = ['order', 'product', 'category', 'price', 'quantity', 'payment', 'customer']
    sales_score = sum(1 for indicator in sales_indicators if any(indicator in col for col in cols_lower))
    
    # Check for campaign indicators
    campaign_indicators = ['cmp', 'campaign', 'accepted', 'response']
    campaign_score = sum(1 for indicator in campaign_indicators if any(indicator in col for col in cols_lower))
    
    if sales_score >= 3:
        return 'SALES'
    elif campaign_score >= 2:
        return 'CAMPAIGN'
    else:
        return 'GENERIC'

def process_sales_data(df):
    """Process e-commerce sales data"""
    metrics = []
    
    # Find key columns
    category_col = next((col for col in df.columns if 'category' in col.lower()), None)
    price_col = next((col for col in df.columns if 'total' in col.lower() and 'price' in col.lower()), None)
    if not price_col:
        price_col = next((col for col in df.columns if 'price' in col.lower()), None)
    quantity_col = next((col for col in df.columns if 'quantity' in col.lower()), None)
    status_col = next((col for col in df.columns if 'status' in col.lower()), None)
    
    if category_col and price_col:
        # Group by category
        category_stats = df.groupby(category_col).agg({
            price_col: ['sum', 'mean', 'count']
        }).reset_index()
        category_stats.columns = ['Category', 'TotalRevenue', 'AvgOrderValue', 'OrderCount']
        
        for _, row in category_stats.iterrows():
            metrics.append({
                'Name': row['Category'],
                'Value': row['TotalRevenue'],
                'Count': row['OrderCount'],
                'Average': row['AvgOrderValue']
            })
    
    return pd.DataFrame(metrics), category_col, price_col

def process_campaign_data(df):
    """Process marketing campaign data"""
    # Find campaign columns
    campaigns = {}
    cmp_columns = [col for col in df.columns if 'cmp' in col.lower() or 'campaign' in col.lower()]
    
    for col in cmp_columns:
        if df[col].dtype in ['int64', 'float64']:
            unique_vals = df[col].dropna().unique()
            if set(unique_vals).issubset({0, 1, 0.0, 1.0}):
                campaigns[col] = col
    
    if 'Response' in df.columns:
        campaigns['Response Campaign'] = 'Response'
    
    # Calculate spending
    spending_patterns = ['mnt', 'amount', 'spend', 'revenue', 'sales']
    spending_cols = [col for col in df.columns if any(p in col.lower() for p in spending_patterns) and df[col].dtype in ['int64', 'float64']]
    
    if spending_cols:
        df['TotalSpend'] = df[spending_cols].sum(axis=1)
    else:
        df['TotalSpend'] = 0
    
    metrics = []
    total_customers = len(df)
    
    for campaign_name, column in campaigns.items():
        acceptances = int(df[column].sum())
        acceptance_rate = (acceptances / total_customers) * 100
        acceptors = df[df[column] == 1]
        avg_spend = acceptors['TotalSpend'].mean() if len(acceptors) > 0 else 0
        
        metrics.append({
            'Name': campaign_name,
            'Value': acceptance_rate,
            'Count': acceptances,
            'Average': avg_spend
        })
    
    return pd.DataFrame(metrics), None, None

# Header
st.markdown('<div class="main-header">🌟 Universal Insight Engine</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-header">Upload ANY CSV - Campaigns, Sales, Analytics - Get Instant Insights</div>', unsafe_allow_html=True)

# Sidebar
with st.sidebar:
    st.header("🎯 Supported Datasets")
    st.success("""
    **📊 Marketing Campaigns**
    - Campaign acceptances
    - Customer spending
    - Performance metrics
    
    **🛒 E-commerce Sales**
    - Product categories
    - Order values
    - Revenue analysis
    
    **📈 Generic Data**
    - Any CSV with metrics
    - Auto-adapts to your schema
    """)
    
    st.header("✨ Smart Features")
    st.info("""
    ✅ Auto-detects data type
    ✅ Adapts charts & insights
    ✅ Handles ANY format
    ✅ Professional reports
    """)

# Main content
uploaded_file = st.file_uploader("📁 Upload your CSV file", type=['csv', 'txt'])

if uploaded_file is not None:
    try:
        # Auto-detect delimiter
        file_content = uploaded_file.read()
        delimiter = detect_delimiter(file_content)
        uploaded_file.seek(0)
        
        # Read CSV
        df = pd.read_csv(uploaded_file, sep=delimiter, encoding='utf-8', on_bad_lines='skip')
        
        # Detect dataset type
        dataset_type = detect_dataset_type(df)
        
        # Display badge
        badge_colors = {'SALES': '#2ca02c', 'CAMPAIGN': '#1f77b4', 'GENERIC': '#ff7f0e'}
        badge_icons = {'SALES': '🛒', 'CAMPAIGN': '📊', 'GENERIC': '📈'}
        
        st.markdown(f"""
        <div class="dataset-badge" style="background-color: {badge_colors[dataset_type]}; color: white;">
            {badge_icons[dataset_type]} Detected: {dataset_type} Dataset
        </div>
        """, unsafe_allow_html=True)
        
        st.success(f"✅ Loaded {len(df):,} rows × {len(df.columns)} columns")
        
        # Show preview
        with st.expander("🔍 Preview Data"):
            st.dataframe(df.head())
        
        # Process button
        if st.button("🚀 Generate Report", type="primary", use_container_width=True):
            with st.spinner(f"🔄 Processing {dataset_type} data..."):
                
                temp_dir = tempfile.mkdtemp()
                chart1_path = os.path.join(temp_dir, "chart1.png")
                chart2_path = os.path.join(temp_dir, "chart2.png")
                
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                status_text.text("📊 Processing data...")
                progress_bar.progress(20)
                
                # Process based on type
                if dataset_type == 'SALES':
                    metrics_df, category_col, price_col = process_sales_data(df)
                    chart1_label = "Revenue by Category"
                    chart2_label = "Order Count by Category"
                    metric1_col = 'Value'
                    metric2_col = 'Count'
                    ylabel1 = "Total Revenue ($)"
                    ylabel2 = "Number of Orders"
                else:  # CAMPAIGN or GENERIC
                    metrics_df, _, _ = process_campaign_data(df)
                    chart1_label = "Performance Metrics"
                    chart2_label = "Engagement Counts"
                    metric1_col = 'Value'
                    metric2_col = 'Count'
                    ylabel1 ="Metric Value"
                    ylabel2 = "Count"
                
                progress_bar.progress(40)
                status_text.text("📈 Creating charts...")
                
                # Chart 1
                plt.figure(figsize=(12, 6))
                colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A', '#98D8C8', '#F7DC6F', '#E8DAEF', '#D5F4E6']
                bars = plt.bar(metrics_df['Name'], metrics_df[metric1_col], 
                              color=colors[:len(metrics_df)], edgecolor='black', linewidth=1.2)
                
                for bar in bars:
                    height = bar.get_height()
                    plt.text(bar.get_x() + bar.get_width()/2., height,
                            f'{height:,.1f}', ha='center', va='bottom', fontsize=10, fontweight='bold')
                
                plt.xlabel('Category', fontsize=12, fontweight='bold')
                plt.ylabel(ylabel1, fontsize=12, fontweight='bold')
                plt.title(chart1_label, fontsize=14, fontweight='bold', pad=20)
                plt.xticks(rotation=45, ha='right')
                plt.grid(axis='y', alpha=0.3, linestyle='--')
                plt.tight_layout()
                plt.savefig(chart1_path, dpi=300, bbox_inches='tight')
                plt.close()
                
                progress_bar.progress(60)
                
                # Chart 2
                plt.figure(figsize=(12, 6))
                bars = plt.bar(metrics_df['Name'], metrics_df[metric2_col], 
                              color=colors[:len(metrics_df)], edgecolor='black', linewidth=1.2)
                
                for bar in bars:
                    height = bar.get_height()
                    plt.text(bar.get_x() + bar.get_width()/2., height,
                            f'{int(height):,}', ha='center', va='bottom', fontsize=10, fontweight='bold')
                
                plt.xlabel('Category', fontsize=12, fontweight='bold')
                plt.ylabel(ylabel2, fontsize=12, fontweight='bold')
                plt.title(chart2_label, fontsize=14, fontweight='bold', pad=20)
                plt.xticks(rotation=45, ha='right')
                plt.grid(axis='y', alpha=0.3, linestyle='--')
                plt.tight_layout()
                plt.savefig(chart2_path, dpi=300, bbox_inches='tight')
                plt.close()
                
                progress_bar.progress(70)
                status_text.text("🤖 Generating insights...")
                
                # Insights with detailed analysis and suggestions
                top_performer = metrics_df.loc[metrics_df[metric1_col].idxmax()]
                worst_performer = metrics_df.loc[metrics_df[metric1_col].idxmin()]
                second_best = metrics_df.nlargest(2, metric1_col).iloc[-1] if len(metrics_df) >= 2 else top_performer
                
                if dataset_type == 'SALES':
                    total_revenue = metrics_df['Value'].sum()
                    total_orders = metrics_df['Count'].sum()
                    avg_order_value = metrics_df['Average'].mean()
                    
                    # Calculate performance gaps
                    revenue_gap = top_performer['Value'] / worst_performer['Value'] if worst_performer['Value'] > 0 else 0
                    
                    # Identify trends
                    top_3 = metrics_df.nlargest(3, 'Value')
                    bottom_3 = metrics_df.nsmallest(3, 'Value')
                    
                    # Calculate additional statistics
                    median_revenue = metrics_df['Value'].median()
                    revenue_std = metrics_df['Value'].std()
                    top_3_concentration = (top_3['Value'].sum()/total_revenue*100)
                    
                    insights = f"""EXECUTIVE SUMMARY
• {len(df):,} orders | ${total_revenue:,.2f} revenue | ${avg_order_value:,.2f} avg order
• {len(metrics_df)} categories | Median revenue: ${median_revenue:,.2f} | Std dev: ${revenue_std:,.2f}

📈 TOP CATEGORY: {top_performer['Name']}
• Revenue: ${top_performer['Value']:,.2f} ({(top_performer['Value']/total_revenue*100):.1f}%)
• Orders: {int(top_performer['Count']):,} | Avg: ${top_performer['Average']:,.2f}
• Analysis: Market leader, {(top_performer['Value']/median_revenue):.1f}x above median

⚠️ UNDERPERFORMER: {worst_performer['Name']}
• Revenue: ${worst_performer['Value']:,.2f} ({(worst_performer['Value']/total_revenue*100):.1f}%)
• Gap: {revenue_gap:.1f}x behind leader | {(median_revenue/worst_performer['Value'] if worst_performer['Value'] > 0 else 0):.1f}x below median

DATA INSIGHTS
• Revenue Distribution: Top 3 = {top_3_concentration:.1f}%, indicating {"high" if top_3_concentration > 60 else "moderate"} concentration
• Performance Spread: {revenue_std/avg_order_value*100:.1f}% variance ({"unstable" if revenue_std/avg_order_value > 0.5 else "stable"} portfolio)
• Growth Potential: {second_best['Name']} at {(second_best['Value']/top_performer['Value']*100):.1f}% of leader
• Market Pattern: {"Winner-take-all" if top_performer['Value']/median_revenue > 3 else "Balanced competition"}

STRATEGIC RECOMMENDATIONS
1. MAXIMIZE WINNERS
   → {top_performer['Name']}: Increase ad spend 20-30% (high ROI expected)
   → Inventory: Stock 1.5x current levels for peak demand
   → Upsell: Bundle with {second_best['Name']} for cross-category growth

2. RESCUE UNDERPERFORMERS
   → {worst_performer['Name']}: Price test (try 15% discount for 2 weeks)
   → Customer surveys: Why low conversion? Fix pain points
   → Placement: Move to high-traffic areas, improve product images

3. PORTFOLIO OPTIMIZATION
   → Risk: {top_3_concentration:.0f}% in 3 categories is {"too high" if top_3_concentration > 70 else "manageable"}
   → Action: Develop categories at 3-7% revenue for stability
   → Target: Achieve <60% concentration in top 3 within 6 months"""

                else:  # CAMPAIGN or GENERIC
                    total_value = metrics_df['Value'].sum() if metric1_col == 'Value' else metrics_df['Count'].sum()
                    avg_performance = metrics_df['Value'].mean()
                    performance_variance = metrics_df['Value'].std()
                    
                    # Calculate efficiency ratio
                    if 'Average' in metrics_df.columns and metrics_df['Average'].sum() > 0:
                        roi_campaigns = metrics_df.nlargest(3, 'Average')
                        efficiency_metric = f"Top ROI: {roi_campaigns.iloc[0]['Name']} (${roi_campaigns.iloc[0]['Average']:.0f}/customer)"
                    else:
                        efficiency_metric = f"Performance Variance: ±{performance_variance:.1f}"
                    
                    # Calculate additional statistics
                    median_performance = metrics_df['Value'].median()
                    q1 = metrics_df['Value'].quantile(0.25)
                    q3 = metrics_df['Value'].quantile(0.75)
                    
                    insights = f"""EXECUTIVE SUMMARY
• {len(df):,} records | Range: {metrics_df['Value'].min():.1f}-{metrics_df['Value'].max():.1f}
• Avg: {avg_performance:.1f} | Median: {median_performance:.1f} | Std Dev: {performance_variance:.1f}

🏆 TOP PERFORMER: {top_performer['Name']}
• Score: {top_performer['Value']:.1f} ({(top_performer['Value']/total_value*100) if total_value > 0 else 0:.1f}%)
• vs Average: {((top_performer['Value'] - avg_performance) / avg_performance * 100) if avg_performance > 0 else 0:+.1f}% | vs Median: {((top_performer['Value']/median_performance - 1)*100) if median_performance > 0 else 0:+.1f}%
• Status: {int(top_performer['Count']):,} count, clearly outperforming

❌ NEEDS ATTENTION: {worst_performer['Name']}
• Score: {worst_performer['Value']:.1f} | Gap: {((avg_performance - worst_performer['Value']) / avg_performance * 100) if avg_performance > 0 else 0:.1f}% below avg
• Quartile: {"Bottom 25%" if worst_performer['Value'] < q1 else "Lower half"} performance tier

DATA INSIGHTS
• Distribution: Q1={q1:.1f}, Median={median_performance:.1f}, Q3={q3:.1f}
• Variance: ±{performance_variance:.1f} ({"high" if performance_variance > avg_performance * 0.3 else "moderate"} volatility)
• {efficiency_metric}
• Pattern: {len(metrics_df[metrics_df['Value'] > avg_performance])} above avg, {len(metrics_df[metrics_df['Value'] < avg_performance])} below

STRATEGIC RECOMMENDATIONS
1. REPLICATE TOP PERFORMANCE
   → Deep dive: Interview {top_performer['Name']} team for success factors
   → Document: Create playbook of winning tactics and strategies
   → Scale: Apply to bottom 50% metrics, target +{((top_performer['Value'] - avg_performance)/2):.1f} improvement

2. IMMEDIATE FIXES (Week 1-2)
   → {worst_performer['Name']}: Emergency review, reallocate 30% budget
   → A/B test: Try 3 different approaches, measure daily
   → Quick wins: Copy tactics from {second_best['Name']} (close to top)

3. SYSTEMATIC OPTIMIZATION (Month 1-3)
   → Target all below-median metrics for 15-20% boost
   → Build dashboards: Real-time monitoring with weekly reviews
   → Set goals: Move bottom quartile to median within 90 days"""

                
                progress_bar.progress(80)
                status_text.text("📑 Creating PowerPoint...")
                
                # PowerPoint
                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                
                # Slide 1
                slide1 = prs.slides.add_slide(prs.slide_layouts[6])
                title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = f"{dataset_type} Performance Report"
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.alignment = 1
                
                date_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
                date_frame = date_box.text_frame
                date_frame.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
                date_para = date_frame.paragraphs[0]
                date_para.font.size = Pt(20)
                date_para.alignment = 1
                
                # Slide 2: First Chart
                slide2 = prs.slides.add_slide(prs.slide_layouts[6])
                title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
                title_frame = title_box.text_frame
                title_frame.text = chart1_label
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(32)
                title_para.font.bold = True
                slide2.shapes.add_picture(chart1_path, Inches(0.75), Inches(1.2), width=Inches(8.5))
                
                # Slide 3: Second Chart
                slide3 = prs.slides.add_slide(prs.slide_layouts[6])
                title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
                title_frame = title_box.text_frame
                title_frame.text = chart2_label
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(32)
                title_para.font.bold = True
                slide3.shapes.add_picture(chart2_path, Inches(0.75), Inches(1.2), width=Inches(8.5))
                
                # Slide 4: Insights & Recommendations (Full Width)
                slide4 = prs.slides.add_slide(prs.slide_layouts[6])
                title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
                title_frame = title_box.text_frame
                title_frame.text = "📊 Insights & Strategic Recommendations"
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                
                # Full-width text box for insights with proper sizing
                insights_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(6))
                insights_frame = insights_box.text_frame
                insights_frame.text = insights
                insights_frame.word_wrap = True
                
                # Smaller font size and minimal spacing
                for paragraph in insights_frame.paragraphs:
                    paragraph.font.size = Pt(9)
                    paragraph.font.name = 'Calibri'
                    paragraph.space_after = Pt(2)
                    paragraph.line_spacing = 1.0
                
                
                # Save PPT
                ppt_buffer = BytesIO()
                prs.save(ppt_buffer)
                ppt_buffer.seek(0)
                
                progress_bar.progress(100)
                status_text.text("✅ Report generated!")
                
                # Display
                st.markdown("---")
                st.subheader("📊 Generated Charts")
                
                col1, col2 = st.columns(2)
                with col1:
                    st.image(chart1_path, use_container_width=True)
                with col2:
                    st.image(chart2_path, use_container_width=True)
                
                st.markdown("---")
                st.subheader("📈 Key Metrics")
                st.dataframe(metrics_df, use_container_width=True)
                
                # Download
                st.markdown("---")
                st.markdown('<div class="success-box">', unsafe_allow_html=True)
                st.markdown("### ✅ Report Ready!")
                st.download_button(
                    label="📥 Download PowerPoint Report",
                    data=ppt_buffer,
                    file_name=f"{dataset_type}_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="primary",
                    use_container_width=True
                )
                st.markdown('</div>', unsafe_allow_html=True)
                
                # Cleanup
                try:
                    os.remove(chart1_path)
                    os.remove(chart2_path)
                    os.rmdir(temp_dir)
                except:
                    pass
                
    except Exception as e:
        st.error(f"❌ Error: {str(e)}")
        st.exception(e)

else:
    st.markdown("---")
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown('<div class="insight-box">', unsafe_allow_html=True)
        st.markdown("### 🛒 E-commerce Sales")
        st.markdown("Orders, products, categories, revenue")
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col2:
        st.markdown('<div class="insight-box">', unsafe_allow_html=True)
        st.markdown("### 📊 Marketing Campaigns")
        st.markdown("Acceptances, spending, performance")
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col3:
        st.markdown('<div class="insight-box">', unsafe_allow_html=True)
        st.markdown("### 📈 Any CSV Data")
        st.markdown("Auto-adapts to your schema")
        st.markdown('</div>', unsafe_allow_html=True)

st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #666;">
    <p>🌟 Universal Engine - Handles ANY CSV Type | Auto-Adapts to Your Data</p>
</div>
""", unsafe_allow_html=True)
